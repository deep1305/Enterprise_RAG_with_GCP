import logfire
from unstructured.partition.auto import partition
from docx import Document
from pptx import Presentation


def _parse_docx_fallback(file_path: str) -> str:
    """
    Lightweight DOCX text extraction fallback.
    """
    doc = Document(file_path)
    parts = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            row_values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_values:
                parts.append(" | ".join(row_values))

    return "\n".join(parts)


def _parse_pptx_fallback(file_path: str) -> str:
    """
    Lightweight PPTX text extraction fallback.
    """
    presentation = Presentation(file_path)
    parts = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"[Slide {slide_index}]")
        for shape in slide.shapes:
            text = ""
            if hasattr(shape, "text"):
                text = shape.text.strip()
            elif hasattr(shape, "text_frame") and shape.text_frame:
                text = shape.text_frame.text.strip()

            if text:
                parts.append(text)

    return "\n".join(parts)

def parse_office(file_path: str):
    """
    Parses Office documents (.docx, .pptx) using the Unstructured library.
    Unlike PDFs, these formats are structured and lightweight, so they are processed locally.
    """
    with logfire.span("📄 Office Document Parsing", filename=file_path):
        extension = file_path.lower().split(".")[-1]

        try:
            # Unstructured automatically detects if it's docx or pptx
            elements = partition(filename=file_path)
            full_text = "\n".join([str(el) for el in elements])
            
            if not full_text.strip():
                logfire.warning(f"⚠️ Unstructured returned empty text for {file_path}")
            else:
                logfire.info(f"✅ Successfully parsed {len(full_text)} characters with unstructured")
                return full_text

        except Exception as unstructured_error:
            logfire.warning(f"⚠️ Unstructured office parse failed, switching to fallback: {unstructured_error}")

        try:
            if extension == "docx":
                full_text = _parse_docx_fallback(file_path)
            elif extension == "pptx":
                full_text = _parse_pptx_fallback(file_path)
            else:
                raise ValueError(f"Unsupported Office extension for fallback: .{extension}")

            if not full_text.strip():
                logfire.warning(f"⚠️ Fallback parser returned empty text for {file_path}")
            else:
                logfire.info(f"✅ Successfully parsed {len(full_text)} characters with fallback parser")

            return full_text
        except Exception as fallback_error:
            logfire.error(f"❌ Office Parse Failed: {fallback_error}")
            raise fallback_error
